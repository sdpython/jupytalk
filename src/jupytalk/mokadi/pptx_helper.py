"""
@file
@brief Helpers on Powerpoint presentation.
It relies on module `python-pptx <https://python-pptx.readthedocs.io/en/latest/>`_.
"""
from pptx.util import Pt
from pptx.dml.color import RGBColor


def pptx_enumerate_text(ptx):
    """
    Enumerate all text content in a presentation.

    @param      ptx         presentation
    @return                 enumerate (slide, shape, zone, text)
    """
    for islide, slide in enumerate(ptx.slides):
        for ishape, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            text_frame = shape.text_frame
            for ip, p in enumerate(text_frame.paragraphs):
                if p.text:
                    yield (islide, ishape, ip, p.text)


def pptx_apply_transform(ptx, replace_func):
    """
    Apply the same transformation on all text zone.

    @param      ptx             presentation
    @param      replace_func    function ``f(islide, ishape, ip, text) --> text
    @return                     number of modified zone

    @warning Updated text does not retain style (bug).
    """
    for islide, slide in enumerate(ptx.slides):
        for ishape, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            text_frame = shape.text_frame
            for ip, p in enumerate(text_frame.paragraphs):
                if p.text:
                    new_text = replace_func(p.text)
                    if new_text != p.text:
                        font = p.font
                        p.text = ""  # new_text
                        run = p.add_run()
                        run.text = new_text
                        run.font.name = font.name
                        run.font.size = Pt(25)  # font.size
                        #run.font.color.rgb = font.color.rgb
                        run.font.color.rgb = RGBColor(0xFF, 0x7F, 0x50)
